#!/usr/bin/env python3
"""Build the GrassApp PowerPoint template (16:9).

Applies the two grassroots backgrounds (tool/make_ppt_bg.py) to the slide
layouts and bakes in the colour scheme:

  * Title Slide layout  -> sunrise background, WHITE title + subtitle.
  * Content layouts      -> gentle blue-wash background, DEEP-BLUE headings,
                            BLACK body text.

Emits two files in assets/:
  * GrassApp.potx          — the reusable template (no slides). Double-click to
                             start a new deck from it; the template is untouched.
  * GrassApp-template.pptx — an example deck (a title + a content slide) so you
                             can see the layouts in use.

Run from glp_multiagent (after make_ppt_bg.py):
    python3 tool/make_ppt_template.py
"""
import os
import zipfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TITLE_BG = "assets/grassapp_ppt_title.png"
CONTENT_BG = "assets/grassapp_ppt_content.png"

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x15, 0x47, 0x7A)     # deep-blue headings
BLACK = RGBColor(0x11, 0x11, 0x11)    # body text

TITLE_LAYOUT = 0          # Title Slide
CONTENT_LAYOUTS = [1, 5, 6]   # Title and Content, Title Only, Blank


def set_bg(layout, path, W, H):
    """Put a full-bleed picture behind everything on a layout.

    LayoutShapes has no add_picture, so go through the part (to register the
    image + relationship) and the underlying spTree's add_pic directly.
    """
    image_part, rId = layout.part.get_or_add_image_part(path)
    tree = layout.shapes._spTree
    sid = max((s.shape_id for s in layout.shapes), default=1) + 1
    pic = tree.add_pic(sid, "Background", "", rId, 0, 0, int(W), int(H))
    tree.remove(pic)
    tree.insert(2, pic)       # behind all placeholders (after nvGrpSpPr, grpSpPr)


def colour_placeholders(layout, title_col, body_col):
    for ph in layout.placeholders:
        idx = ph.placeholder_format.idx
        col = title_col if idx == 0 else body_col
        for para in ph.text_frame.paragraphs:
            if para.font is not None:
                para.font.color.rgb = col
                para.font.name = "Arial"


def style(run_or_para, col, size=None, bold=None):
    f = run_or_para.font
    f.color.rgb = col
    f.name = "Arial"
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold


def base():
    """A 16:9 presentation with backgrounds + colours baked into the layouts."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height
    set_bg(prs.slide_layouts[TITLE_LAYOUT], TITLE_BG, W, H)
    colour_placeholders(prs.slide_layouts[TITLE_LAYOUT], WHITE, WHITE)
    for i in CONTENT_LAYOUTS:
        set_bg(prs.slide_layouts[i], CONTENT_BG, W, H)
        colour_placeholders(prs.slide_layouts[i], BLUE, BLACK)
    return prs


def add_examples(prs):
    s = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
    s.shapes.title.text = "GrassApp"
    style(s.shapes.title.text_frame.paragraphs[0], WHITE, 54, True)
    sub = s.placeholders[1]
    sub.text = "Grassroots digital cooperatives"
    style(sub.text_frame.paragraphs[0], WHITE, 24, False)

    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Three platforms, one app"
    style(s.shapes.title.text_frame.paragraphs[0], BLUE, 40, True)
    body = s.placeholders[1].text_frame
    bullets = ["Friends — the social graph",
               "Coins — coins among friends",
               "Chats — the social network"]
    body.text = bullets[0]
    style(body.paragraphs[0], BLACK, 28, False)
    for b in bullets[1:]:
        p = body.add_paragraph()
        p.text = b
        style(p, BLACK, 28, False)


def to_template(pptx_path, potx_path):
    """Re-stamp a saved .pptx as a PowerPoint template (.potx).

    The only difference is the content-type of the presentation part, so rewrite
    that one override in [Content_Types].xml.
    """
    PRES = ("application/vnd.openxmlformats-officedocument"
            ".presentationml.presentation.main+xml")
    TMPL = ("application/vnd.openxmlformats-officedocument"
            ".presentationml.template.main+xml")
    with zipfile.ZipFile(pptx_path) as zin:
        items = [(i, zin.read(i.filename)) for i in zin.infolist()]
    with zipfile.ZipFile(potx_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for info, data in items:
            if info.filename == "[Content_Types].xml":
                data = data.replace(PRES.encode(), TMPL.encode())
            zout.writestr(info, data)


# --- The reusable template (.potx): layouts only, no slides -----------------
base().save("assets/_tmp_template.pptx")
to_template("assets/_tmp_template.pptx", "assets/GrassApp.potx")
os.remove("assets/_tmp_template.pptx")

# --- The example deck (.pptx): same layouts, with two example slides --------
prs = base()
add_examples(prs)
prs.save("assets/GrassApp-template.pptx")

print("wrote assets/GrassApp.potx and assets/GrassApp-template.pptx")
